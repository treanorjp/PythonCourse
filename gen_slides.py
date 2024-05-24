import os
from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
filename = 'Python_Data_Structures.pptx'

if os.path.exists(filename):
    prs = Presentation(filename)
else:
    prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout for title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Introduction to Python Data Structures"
subtitle.text = "Strings, Lists, Tuples, and Dictionaries"

# Strings Section

# Slide 1: What is a String?
slide_layout = prs.slide_layouts[1]  # 1 is the layout for title and content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "What is a String?"
content.text = """
- A string is a sequence of characters.
- Strings are used to represent text.
- In Python, strings are immutable, meaning they cannot be changed after creation.
- Strings can be enclosed in single quotes (' '), or double quotes (" ").
"""

# Slide 2: Creating Strings
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Creating Strings"
content.text = """
- Single quotes: 'Hello'
- Double quotes: "Hello"
"""

# Slide 3: Accessing Characters
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Accessing Characters"
content.text = """
- Zero-based indexing
- Accessing the first character: my_string[0]
- Accessing the last character: my_string[-1]
"""

# Slide 4: Slicing Strings
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Slicing Strings"
content.text = """
- Access a substring: my_string[1:4]
- Slicing with a step: my_string[0:5:2]
"""

# Slide 5: String Methods
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "String Methods"
content.text = """
- Convert to uppercase: my_string.upper()
- Convert to lowercase: my_string.lower()
- Split a string: my_string.split()
- Join a list of strings: ' '.join(list_of_strings)
"""

# Slide 6: String Concatenation
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "String Concatenation"
content.text = """
- Using + operator: 'Hello' + ' ' + 'World'
- Using join method: ' '.join(['Hello', 'World'])
"""

# Slide 7: String Formatting
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "String Formatting"
content.text = """
- Using str.format(): 'Hello {}'.format('World')
- Using f-strings (Python 3.6+): f'Hello {name}'
- Limiting a float to n decimal places: '{:.2f}'.format(3.20159)
"""

# Slide 8: Common String Operations
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Common String Operations"
content.text = """
- Find a substring: my_string.find('substring')
- Replace a substring: my_string.replace('old', 'new')
- Check if alphanumeric: my_string.isalnum()
- Strip whitespace: my_string.strip()
"""

# Lists Section

# Slide 9: What is a List?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "What is a List?"
content.text = """
- A list is an ordered collection of items.
- Lists are mutable, meaning they can be modified.
- Lists can contain elements of different data types.
- Created using square brackets: [1, 2, 3]
"""

# Slide 10: Creating a List
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Creating a List"
content.text = """
- An empty list: empty_list = []
- A list of integers: int_list = [1, 2, 3, 4, 5]
- A list of strings: str_list = ["apple", "banana", "cherry"]
- A mixed list: mixed_list = [1, "apple", 3.20, True]
"""

# Slide 11: Accessing Elements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Accessing Elements"
content.text = """
- Zero-based indexing
- Accessing the first element: my_list[0]
- Accessing the last element: my_list[-1]
"""

# Slide 12: Modifying Elements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Modifying Elements"
content.text = """
- Modify an element by assigning a new value: my_list[1] = "blueberry"
"""

# Slide 13: Adding Elements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Adding Elements"
content.text = """
- Append: my_list.append("cherry")
- Insert: my_list.insert(1, "blueberry")
- Extend: my_list.extend(["date", "elderberry"])
"""

# Slide 14: Removing Elements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Removing Elements"
content.text = """
- Remove by value: my_list.remove("banana")
- Remove by index and return: my_list.pop(0)
- Delete by index: del my_list[0]
"""

# Slide 15: Slicing
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Slicing"
content.text = """
- Access a subset of a list: my_list[1:3]
- Slicing with a step: my_list[0:4:2]
"""

# Slide 16: List Comprehensions
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "List Comprehensions"
content.text = """
- Create lists concisely
- Example: squares = [x**2 for x in range(5)]
"""

# Slide 17: Common List Methods
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Common List Methods"
content.text = """
- len(list)
- list.sort()
- list.reverse()
- list.index(value)
- list.count(value)
"""

# Tuples Section

# Slide 18: What is a Tuple?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "What is a Tuple?"
content.text = """
- A tuple is an ordered collection of items.
- Tuples are immutable, meaning they cannot be modified after creation.
- Tuples can contain elements of different data types.
- Created using parentheses: (1, 2, 3)
"""

# Slide 19: Creating a Tuple
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Creating a Tuple"
content.text = """
- An empty tuple: empty_tuple = ()
- A tuple of integers: int_tuple = (1, 2, 3, 4, 5)
- A tuple of strings: str_tuple = ("apple", "banana", "cherry")
- A mixed tuple: mixed_tuple = (1, "apple", 3.14, True)
"""

# Slide 20: Accessing Elements in a Tuple
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Accessing Elements in a Tuple"
content.text = """
- Zero-based indexing
- Accessing the first element: my_tuple[0]
- Accessing the last element: my_tuple[-1]
"""

# Slide 21: Common Tuple Operations
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Common Tuple Operations"
content.text = """
- len(tuple)
- tuple.index(value)
- tuple.count(value)
"""

# Dictionaries Section

# Slide 22: What is a Dictionary?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "What is a Dictionary?"
content.text = """
- A dictionary is an unordered collection of key-value pairs.
- Keys must be unique and immutable (e.g., strings, numbers, tuples).
- Values can be of any data type and can be duplicated.
- Created using curly braces: {'key1': 'value1', 'key2': 'value2'}
"""

# Slide 23: Creating a Dictionary
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Creating a Dictionary"
content.text = """
- An empty dictionary: empty_dict = {}
- A dictionary with key-value pairs: my_dict = {'name': 'Alice', 'age': 25}
"""

# Slide 24: Accessing Values
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Accessing Values"
content.text = """
- Access value by key: my_dict['name']
- Using get() method: my_dict.get('name')
"""

# Slide 25: Modifying Values
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Modifying Values"
content.text = """
- Change an existing value: my_dict['age'] = 26
- Add a new key-value pair: my_dict['city'] = 'New York'
"""

# Slide 26: Removing Key-Value Pairs
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Removing Key-Value Pairs"
content.text = """
- Using pop() method: my_dict.pop('age')
- Using del statement: del my_dict['name']
- Using popitem() method to remove the last inserted item: my_dict.popitem()
"""

# Slide 27: Dictionary Methods
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Dictionary Methods"
content.text = """
- keys(): Returns a view object of all keys
- values(): Returns a view object of all values
- items(): Returns a view object of all key-value pairs
- update(): Updates the dictionary with elements from another dictionary
"""

# Slide 28: Looping Through a Dictionary
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Looping Through a Dictionary"
content.text = """
- Loop through keys: for key in my_dict
- Loop through values: for value in my_dict.values()
- Loop through key-value pairs: for key, value in my_dict.items()
"""

# Slide 29: Dictionary Comprehensions
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Dictionary Comprehensions"
content.text = """
- Create dictionaries concisely
- Example: squares = {x: x**2 for x in range(5)}
"""

# Slide 30: Nested Dictionaries
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Nested Dictionaries"
content.text = """
- A dictionary within a dictionary
- Example: nested_dict = {'parent': {'child': 'value'}}
"""

# Slide 31: Differences Between Lists and Dictionaries
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Differences Between Lists and Dictionaries"
content.text = """
- Lists are ordered collections, while dictionaries are unordered.
- Lists use integer indices to access elements, while dictionaries use keys.
- Lists are mutable and allow duplicate elements, while dictionaries require unique keys.
- Use lists for ordered data and dictionaries for key-value pairs.
"""

# Save the presentation
prs.save('Python_Data_Structures.pptx')