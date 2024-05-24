from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation('OOP_in_Python.pptx')

# Title Slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout for title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Overview of Object-Oriented Programming in Python"
subtitle.text = "Understanding OOP Concepts in Python"

# Slide 1: What is Object-Oriented Programming?
slide_layout = prs.slide_layouts[1]  # 1 is the layout for title and content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "What is Object-Oriented Programming?"
content.text = """
- Object-Oriented Programming (OOP) is a programming paradigm based on the concept of objects.
- Objects are instances of classes, which can contain data and methods.
- OOP aims to implement real-world entities like inheritance, polymorphism, encapsulation, and abstraction.
"""

# Slide 2: Classes and Objects
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Classes and Objects"
content.text = """
- A class is a blueprint for creating objects.
- An object is an instance of a class.
- Example:
  class Dog:
      def __init__(self, name, age):
          self.name = name
          self.age = age
  my_dog = Dog("Buddy", 3)
"""

# Slide 3: Attributes and Methods
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Attributes and Methods"
content.text = """
- Attributes are variables that belong to an object.
- Methods are functions that belong to an object.
- Example:
  class Dog:
      def __init__(self, name, age):
          self.name = name  # Attribute
          self.age = age    # Attribute
      def bark(self):       # Method
          print("Woof!")
"""

# Slide 4: Inheritance
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Inheritance"
content.text = """
- Inheritance allows a class to inherit attributes and methods from another class.
- The class that inherits is called a subclass, and the class being inherited from is called a superclass.
- Example:
  class Animal:
      def __init__(self, name):
          self.name = name
  class Dog(Animal):
      def bark(self):
          print("Woof!")
  my_dog = Dog("Buddy")
  my_dog.bark()
"""

# Slide 5: Polymorphism
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Polymorphism"
content.text = """
- Polymorphism allows methods to do different things based on the object it is acting upon.
- Example:
  class Dog:
      def speak(self):
          return "Woof!"
  class Cat:
      def speak(self):
          return "Meow!"
  animals = [Dog(), Cat()]
  for animal in animals:
      print(animal.speak())
"""

# Slide 6: Encapsulation
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Encapsulation"
content.text = """
- Encapsulation is the concept of wrapping data and methods into a single unit (class).
- It restricts direct access to some of the object's components.
- Example:
  class Person:
      def __init__(self, name, age):
          self.__name = name  # Private attribute
          self.__age = age    # Private attribute
      def get_name(self):
          return self.__name
  person = Person("Alice", 30)
  print(person.get_name())
"""

# Slide 7: Abstraction
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Abstraction"
content.text = """
- Abstraction is the concept of hiding the complex implementation details and showing only the necessary features.
- It can be achieved using abstract classes and interfaces.
- Example:
  from abc import ABC, abstractmethod
  class Animal(ABC):
      @abstractmethod
      def make_sound(self):
          pass
  class Dog(Animal):
      def make_sound(self):
          return "Woof!"
  dog = Dog()
  print(dog.make_sound())
"""

# Slide 8: Benefits of OOP
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Benefits of OOP"
content.text = """
- Modularity: Code is organized into objects.
- Reusability: Classes can be reused across programs.
- Extensibility: New functionality can be added with minimal changes to existing code.
- Maintainability: Code is easier to manage and maintain.
"""

# Slide 9: OOP in Python
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "OOP in Python"
content.text = """
- Python is an object-oriented programming language.
- It supports all the key concepts of OOP: classes, objects, inheritance, polymorphism, encapsulation, and abstraction.
- Python's simplicity and readability make it a great choice for learning and implementing OOP.
"""

# Save the presentation
prs.save('OOP_in_Python.pptx')